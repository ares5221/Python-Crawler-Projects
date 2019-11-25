#!/usr/bin/env python
# _*_ coding:utf-8 _*_
import os
from docx import Document
from win32com import client as wc
# pip install python-pptx
from pptx import Presentation
# pip install wand
from wand.image import Image
# pip install baidu-aip
from aip import AipOcr
import zipfile
import rarfile

# from unrar import rarfile

OS_Platform = 'win10'


def text_content_analysis(key):
    create_classifer_files(key)
    base_dir = './../data_v2/source'
    curr_dir = base_dir + '/' + key
    for file_item in os.listdir(curr_dir):
        curr_file_name = os.path.join(curr_dir, file_item)
        # print(curr_file_name)
        # todo
        # https://www.jianshu.com/p/056e94ca301e
        if file_item[-5:] == '.docx':
            phrase_docx(curr_file_name)
        if file_item[-4:] == '.doc':
            phrase_doc(curr_file_name)
        if file_item[-4:] == '.txt':
            phrase_txt(curr_file_name)
        if file_item[-4:] == '.pdf':
            phrase_pdf(curr_file_name)
        if file_item[-4:] == '.ppt':
            phrase_ppt(curr_file_name)
        if file_item[-5:] == '.ppxt':
            readPptx(curr_file_name)
        if file_item[-4:] == '.zip':
            phrase_zip(curr_file_name)
        if file_item[-4:] == '.rar':
            phrase_rar(curr_file_name)


def create_classifer_files(key):
    print('创建分类文件夹，若已经存在则跳过')
    base_dir = './../data_v2/target/'
    curr_dir = os.path.join(base_dir, key)
    if not os.path.exists(curr_dir):
        os.makedirs(curr_dir)
    if not os.path.exists(os.path.join(curr_dir, '试题')):
        os.makedirs(os.path.join(curr_dir, '试题'))
    if not os.path.exists(os.path.join(curr_dir, '资料')):
        os.makedirs(os.path.join(curr_dir, '资料'))
    if not os.path.exists(os.path.join(curr_dir, '课件')):
        os.makedirs(os.path.join(curr_dir, '课件'))
    if not os.path.exists(os.path.join(curr_dir, '其他')):
        os.makedirs(os.path.join(curr_dir, '其他'))


def phrase_doc(file_path):
    abs_base_dir = os.path.abspath('./../')
    sp = os.path.split(file_path)
    print('通过文件名和文件内容来分析文件类型, 当前路径', file_path)
    curr_key_word_file_path = sp[0]
    file_name = sp[1]  # 获取文件名

    sp1 = curr_key_word_file_path.split('/')  # 切割路径，拼接绝对路径
    key_word = sp1[4]  # 根据切割得到文件所属关键字类别
    if OS_Platform == 'win10':
        abs_win10_path = os.path.join(abs_base_dir, sp1[2], sp1[3], sp1[4], file_name)
        final_path = abs_win10_path
    else:
        pass  # todo 设置linux的路径
    content = []  # 存储文件内容
    # 读取doc文件内容
    word = wc.Dispatch('Word.Application')
    # win32 读取doc文件，无法识别相对路径，只能识别绝对路径 对于linux需要重新设置
    doc = word.Documents.Open(final_path)
    for para in doc.paragraphs:
        content.append(para.Range.Text)
    doc.Close()
    # 根据文件名及文件内容分析文件具体类型
    class_type = get_content_type(file_name, content)

    print(file_path, class_type)
    save_corresponding_path(file_path, class_type, key_word)


def save_corresponding_path(file_path, class_type, key_word):
    base_dir = './../data_v2/target/'
    backup_dir = './../data_v2/backup'
    file_name = os.path.split(file_path)[1]
    save_path = os.path.join(base_dir, key_word, class_type, file_name)
    if not os.path.exists(save_path):
        os.rename(file_path, save_path)
    else:
        bc_path = os.path.join(backup_dir, file_name)
        os.rename(file_path, bc_path)
        print('文件已经存在，将文件移动到backup 做删除')


def get_content_type(file_name, content):
    full_text = ''
    for cc in content:
        full_text += cc
    if '试题' in file_name or ('答题表' in full_text and '班级' in full_text and '姓名' in full_text):
        return '试题'
    elif '练习' in file_name or '同步练习' in full_text or '基础练习' in full_text or '资料' in file_name:
        return '资料'
    elif '课件' in file_name or '笔记' in full_text or '专题' in file_name or '笔记' in file_name:
        return '课件'
    else:
        return '其他'


def phrase_docx(file_path):
    sp1 = file_path.split('/')  # 切割路径，拼接绝对路径
    key_word = os.path.split(os.path.split(file_path)[0])[1]  # 根据切割得到文件所属关键字类别
    # 读取doc文件内容
    document = Document(file_path)  # 打开docx文件
    content = []  # 存储文件内容
    for paragraph in document.paragraphs:
        content.append(paragraph.text)
    class_type = get_content_type(sp1[-1], content)
    save_corresponding_path(file_path, class_type, key_word)


def phrase_txt(file_path):
    sp1 = file_path.split('/')  # 切割路径，拼接绝对路径
    key_word = os.path.split(os.path.split(file_path)[0])[1]  # 根据切割得到文件所属关键字类别
    # 读取txt文件内容
    content = []  # 存储文件内容
    if os.path.exists(file_path):
        with open(file_path, 'r', encoding='utf-8') as f:
            for l in f:
                temp = l.rstrip('\n').rstrip().split('\t')[0]
                content += temp.replace(' ', '')
        class_type = get_content_type(sp1[-1], content)
        save_corresponding_path(file_path, class_type, key_word)


def phrase_pdf(file_path):
    sp1 = file_path.split('/')  # 切割路径，拼接绝对路径
    key_word = os.path.split(os.path.split(file_path)[0])[1]  # 根据切割得到文件所属关键字类别
    # 读取pdf文件内容,将pdf文件转为jpg图片文件
    content = ""
    # 使用 wand 异常，缺少 ImageMagick 支持
    # http://docs.wand-py.org/en/latest/guide/install.html#install-imagemagick-on-windows
    # https://imagemagick.org/script/download.php#windows
    # 使用 wand 异常，FailedToExecuteCommand `"gswin32c.exe"
    # http://ghostscript.com/download/gsdnld.html
    image_pdf = Image(filename=file_path, resolution=300)
    image_jpeg = image_pdf.convert('jpg')

    # wand已经将PDF中所有的独立页面都转成了独立的二进制图像对象。我们可以遍历这个大对象，并把它们加入到req_image序列中去。
    req_image = []
    for img in image_jpeg.sequence:
        img_page = Image(image=img)
        req_image.append(img_page.make_blob('jpg'))
    # 遍历req_image,保存为图片文件
    for img in req_image:
        ff = open(file_path + '.jpg', 'wb')
        ff.write(img)
        ff.close()
        # 调用图片文字识别
        content += readImage(file_path + '.jpg')
        # 移除临时图片
        os.remove(file_path + '.jpg')
    class_type = get_content_type(sp1[-1], [content])
    print('pdf文件转换为图片后读取的内容信息类型为', class_type)
    save_corresponding_path(file_path, class_type, key_word)


# 读取 图片 文件，返回文件内容
def readImage(fileUrl):
    content = ""
    if os.path.exists(fileUrl):
        APP_ID = '17848885'
        API_KEY = 'oKhR50OyDdO5bcpXhTAPHInh'
        SECRET_KEY = 'oGbbyUQysuGKzySySYhBMWrB4plPatvq'
        client = AipOcr(APP_ID, API_KEY, SECRET_KEY)
        with open(fileUrl, 'rb') as f:
            img = f.read()
            msg = client.basicGeneral(img)
            for i in msg.get('words_result'):
                temp = i.get('words')
                content += temp.replace(' ', '')
    return content


def phrase_ppt(file_path):
    sp1 = file_path.split('/')  # 切割路径，拼接绝对路径
    key_word = os.path.split(os.path.split(file_path)[0])[1]  # 根据切割得到文件所属关键字类别
    # 读取 ppt 文件，安装 pypiwin32，操作本地ppt程序，将ppt转为pptx，再调用读取 pptx 文件方法
    extend_table = False
    extend_image = False
    AbsolutePath = os.path.abspath(file_path)
    powerpoint = wc.Dispatch('PowerPoint.Application')
    ppt = powerpoint.Presentations.Open(AbsolutePath)
    # 保存临时文件
    ppt.SaveAs(AbsolutePath + ".pptx")
    powerpoint.Quit()
    content = readPptx(file_path + ".pptx", extend_table, extend_image)
    # 移除临时文件
    os.remove(file_path + ".pptx")
    class_type = get_content_type(sp1[-1], [content])
    save_corresponding_path(file_path, class_type, key_word)


# 读取 pptx 文件 ,默认读取正文，默认不读取表格，默认不读取图片
def readPptx(fileUrl, extend_table=False, extend_image=False):
    sp1 = fileUrl.split('/')  # 切割路径，拼接绝对路径
    key_word = os.path.split(os.path.split(fileUrl)[0])[1]  # 根据切割得到文件所属关键字类别
    content = ""
    ppt = Presentation(fileUrl)
    for slide in ppt.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                # 提取图片文字
                if extend_image and hasattr(shape, 'image'):
                    # 图片存储本地
                    with open(shape.image.filename, 'wb') as f:
                        f.write(shape.image.blob)
                        f.close()
                    # 调用图片文字识别
                    content += readImage(shape.image.filename)
                    # 移除临时图片
                    os.remove(shape.image.filename)
                # 提取表格内容
                if extend_table and shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            content += cell.text
            else:
                content += shape.text
    content = "".join(content.split())
    class_type = get_content_type(sp1[-1], content)
    save_corresponding_path(fileUrl, class_type, key_word)
    return content


def phrase_zip(file_path):
    key_word = os.path.split(os.path.split(file_path)[0])[1]  # 根据切割得到文件所属关键字类别
    print('当前zip文件所属关键字类别：', key_word)

    abs_base_dir = os.path.abspath('./../')
    sp = os.path.split(file_path)
    curr_key_word_file_path = sp[0]
    file_name = sp[1]  # 获取文件名
    sp1 = curr_key_word_file_path.split('/')  # 切割路径，拼接绝对路径
    if OS_Platform == 'win10':
        abs_win10_path = os.path.join(abs_base_dir, sp1[2], sp1[3], sp1[4], file_name)
        final_path = abs_win10_path
    else:
        pass  # todo 设置linux的路径
    # zipfile 也是必须使用绝对路径
    zzz = zipfile.ZipFile(final_path, 'r')
    for f_name in zzz.namelist():  # z.namelist() 会返回压缩包内所有文件名的列表。
        print(f_name)
    zzz.close()
    # todo 当前暂时不对zip文件内容做复杂提取分析，简单将zip文件统一归类为 其他 即可
    # class_type = get_content_type(sp1[-1], [content])
    # print('pdf文件转换为图片后读取的内容信息类型为', class_type)
    class_type = '其他'
    save_corresponding_path(file_path, class_type, key_word)
    return


def phrase_rar(file_path):
    # todo 先将文件解压再处理
    key_word = os.path.split(os.path.split(file_path)[0])[1]  # 根据切割得到文件所属关键字类别
    print('当前zip文件所属关键字类别：', key_word)
    abs_base_dir = os.path.abspath('./../')
    sp = os.path.split(file_path)
    curr_key_word_file_path = sp[0]
    file_name = sp[1]  # 获取文件名
    sp1 = curr_key_word_file_path.split('/')  # 切割路径，拼接绝对路径
    if OS_Platform == 'win10':
        abs_win10_path = os.path.join(abs_base_dir, sp1[2], sp1[3], sp1[4], file_name)
        final_path = abs_win10_path
    else:
        pass  # todo 设置linux的路径

    path2 = "E:\\New"
    print(final_path)
    # rf = rarfile.RarFile(final_path)# 待解压文件
    # rf.extractall(path2)
    # todo 当前暂时不对zip文件内容做复杂提取分析，简单将zip文件统一归类为 其他 即可
    # class_type = get_content_type(sp1[-1], [content])
    # print('pdf文件转换为图片后读取的内容信息类型为', class_type)
    class_type = '其他'
    save_corresponding_path(file_path, class_type, key_word)


'''
def change_doc2docx(file_path):
    word = wc.Dispatch('Word.Application')
    doc = word.Documents.Open(file_path)
    newname = save_dir + '/' + fname[:-4] + '.docx'
    # print(curr_file_name)
    # print(newname)
    # print(fname,'ssssss')
    doc.SaveAs(newname, 16)
    doc.Close()
'''

if __name__ == '__main__':
    key = '光合作用'
    text_content_analysis(key)
